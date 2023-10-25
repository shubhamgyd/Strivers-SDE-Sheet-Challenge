from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

def pdf_to_images(pdf_path):
    images = convert_from_path(pdf_path)
    return images

def create_ppt(images, output_ppt_path):
    prs = Presentation()

    for idx, image in enumerate(images):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use blank slide layout
        pic = slide.shapes.add_picture(image, Inches(0), Inches(0), width=Inches(10), height=Inches(7))

    prs.save(output_ppt_path)

if __name__ == "__main__":
    pdf_path = "C:/Users/shubh/OneDrive/Desktop/DP-Knapsack Problem-2023.pdf"  # Replace with your PDF file path
    output_ppt_path = "DP-Knapsack Problem-2023.pptx"  # Desired output PowerPoint file path

    images = pdf_to_images(pdf_path)
    create_ppt(images, output_ppt_path)
